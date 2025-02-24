from typing import Dict

from langchain.schema import HumanMessage
from pptx import Presentation

from .base_agent import BaseAgent


class ImprovementAgent(BaseAgent):
    def __init__(self):
        super().__init__()
        self.name = "Improvement Agent"
        self.description = "implementing improvements based on review feedback"

    def _create_improvement_prompt(self, slide: Dict, feedback: Dict) -> str:
        """Create prompt for improving a slide based on feedback"""
        return f"""Improve this presentation slide based on the review feedback.
                Focus on technical accuracy, grammar, and timing.

                Current slide:
                Title: {slide['title']}
                Content: {slide['content']}

                Feedback:
                {feedback}

                Provide improved content that addresses the feedback."""

    def _apply_improvements(self, prs: Presentation, slide_index: int,
                            improved_content: Dict) -> None:
        """Apply improvements to a specific slide"""
        if slide_index >= len(prs.slides):
            return

        self.logger.info(f"Applying improvements to slide {slide_index + 1}")
        slide = prs.slides[slide_index]

        # Update title if provided
        if "title" in improved_content and slide.shapes.title:
            slide.shapes.title.text = improved_content["title"]

        # Update content if provided
        if "content" in improved_content:
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Body placeholder
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = improved_content["content"]

    async def process(self, state: Dict) -> Dict:
        """Implement improvements based on review feedback"""
        try:
            self.logger.info("Starting presentation improvements")

            if "review_feedback" not in state or "presentation" not in state:
                raise ValueError(
                    "Missing review feedback or presentation in state")

            feedback = state["review_feedback"]
            presentation = state["presentation"]

            # Load existing presentation
            self.logger.info("Loading existing presentation")
            prs = Presentation("output/presentation.pptx")

            # Track improvements
            improvements_made = []
            improved_slides = []

            # Process each slide
            self.logger.info(
                f"Processing {len(presentation['slides'])} slides")
            for i, slide in enumerate(presentation["slides"]):
                slide_feedback = feedback["slide_specific_feedback"].get(
                    f"Slide {i+1}", [])

                if slide_feedback:
                    self.logger.info(f"Improving slide {i + 1}")
                    # Get improvement suggestions from LLM
                    messages = [
                        HumanMessage(content=self._create_improvement_prompt(
                            slide, "\n".join(slide_feedback)))
                    ]
                    improvement_response = await self.llm.ainvoke(messages)

                    # Apply improvements
                    improved_content = {
                        "title": slide["title"],
                        "content": improvement_response.content
                    }
                    self._apply_improvements(prs, i, improved_content)

                    improvements_made.append(f"Improved Slide {i+1}")
                    improved_slides.append(improved_content)
                else:
                    self.logger.info(
                        f"No improvements needed for slide {i + 1}")
                    improved_slides.append(slide)

            # Save improved presentation
            self.logger.info("Saving improved presentation")
            prs.save("output/presentation_improved.pptx")

            # Update state
            state["presentation"] = {
                "slides": improved_slides,
                "total_time": presentation["total_time"],
                "current_version": presentation["current_version"] + 1
            }
            state["status"] = "Improvements applied"
            state["current_agent"] = "Review Agent"

            self.logger.success(
                f"Successfully improved {len(improvements_made)} slides")

        except Exception as e:
            self.logger.error(f"Error applying improvements: {str(e)}")
            state["errors"].append(f"Improvement Agent error: {str(e)}")
            state["status"] = "Error in improvements"

        return state
