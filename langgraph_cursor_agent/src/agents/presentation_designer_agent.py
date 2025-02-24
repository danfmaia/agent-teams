from typing import Dict

from pptx import Presentation
from pptx.util import Pt

from .base_agent import BaseAgent


class PresentationDesignerAgent(BaseAgent):
    def __init__(self):
        super().__init__()
        self.name = "Presentation Designer Agent"
        self.description = "creating and formatting PowerPoint presentations"

    def _create_title_slide(self, prs: Presentation, title: str) -> None:
        """Create the title slide"""
        self.logger.info("Creating title slide")
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)

    def _create_section_slide(self, prs: Presentation, section: Dict[str, str]) -> None:
        """Create a slide for a section"""
        self.logger.info(f"Creating slide for section: {section['title']}")
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)

        # Set title
        title_placeholder = slide.shapes.title
        title_placeholder.text = section["title"]
        title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)

        # Set content
        content = slide.placeholders[1]
        content_frame = content.text_frame

        # Process content into bullet points
        lines = section["content"].split("\n")
        current_level = 0
        for line in lines:
            if line.strip():
                p = content_frame.add_paragraph()
                p.text = line.strip()
                p.level = current_level
                p.font.size = Pt(24)

    def _estimate_presentation_time(self, num_slides: int) -> int:
        """Estimate presentation time in minutes"""
        # Basic estimation: 1-2 minutes per slide
        return max(5, num_slides * 1.5)

    async def process(self, state: Dict) -> Dict:
        """Create PowerPoint presentation from translated content"""
        try:
            self.logger.info("Starting presentation creation")

            if "translated_content" not in state:
                raise ValueError("No translated content found in state")

            translated_content = state["translated_content"]

            # Create presentation
            self.logger.info("Initializing PowerPoint presentation")
            prs = Presentation()

            # Set slide size to 16:9
            prs.slide_width = 12192000  # 33.867 cm
            prs.slide_height = 6858000  # 19.05 cm

            # Create title slide
            self._create_title_slide(prs, "Throughput Estimation")

            # Create content slides
            self.logger.info("Creating content slides")
            for section in translated_content["sections"]:
                self._create_section_slide(prs, section)

            # Save presentation
            self.logger.info("Saving presentation")
            output_path = "output/presentation.pptx"
            prs.save(output_path)

            # Update state with slides and timing info
            slides = [{"title": s["title"], "content": s["content"]}
                      for s in translated_content["sections"]]
            total_time = self._estimate_presentation_time(len(slides) + 1)

            state["presentation"] = {
                "slides": slides,
                "total_time": total_time,
                "current_version": 1
            }

            state["status"] = "Presentation created"
            state["current_agent"] = "Review Agent"

            self.logger.success(
                f"Successfully created presentation with {len(slides)} slides")

        except Exception as e:
            self.logger.error(f"Error creating presentation: {str(e)}")
            state["errors"].append(
                f"Presentation Designer Agent error: {str(e)}")
            state["status"] = "Error in presentation design"

        return state
