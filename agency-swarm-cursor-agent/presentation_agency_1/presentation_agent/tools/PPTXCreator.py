from pydantic import Field
import os
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import io

load_dotenv()


class PPTXCreator:
    """
    A tool for creating PowerPoint presentations in Portuguese with text, formulas, and graphs.
    """

    def __init__(self, content, output_path):
        self.content = content
        self.output_path = output_path

    def run(self):
        """
        Creates a PowerPoint presentation with the provided content.
        """
        try:
            # Get the absolute path for the output file
            base_dir = os.path.dirname(os.path.dirname(
                os.path.dirname(os.path.abspath(__file__))))
            full_output_path = os.path.join(base_dir, self.output_path)

            # Ensure output directory exists
            os.makedirs(os.path.dirname(full_output_path), exist_ok=True)

            print(f"Creating presentation at: {full_output_path}")

            prs = Presentation()

            # Set slide width and height for 16:9 aspect ratio
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)

            # Title slide
            self._add_title_slide(prs)

            # Introduction slide
            self._add_introduction_slide(prs)

            # Add content slides for each section
            sections = self.content['sections']

            # Abstract slide
            if sections['abstract']:
                self._add_section_slide(
                    prs, "Resumo", self._summarize_text(sections['abstract'], 200))

            # Methodology slide
            if sections['methodology']:
                self._add_section_slide(
                    prs, "Metodologia", self._summarize_text(sections['methodology'], 200))

            # Results slide
            if sections['results']:
                self._add_section_slide(
                    prs, "Resultados", self._summarize_text(sections['results'], 200))

            # Formulas slide
            if self.content['formulas']:
                self._add_formulas_slide(prs)

            # Conclusion slide
            if sections['conclusion']:
                self._add_section_slide(
                    prs, "Conclusão", self._summarize_text(sections['conclusion'], 200))

            # Save the presentation
            prs.save(full_output_path)
            return f"Presentation created successfully at {full_output_path}"

        except Exception as e:
            return f"Error creating presentation: {str(e)}"

    def _add_title_slide(self, prs):
        """Add the title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Estimativa de Throughput"
        subtitle.text = "Análise de Sistemas Computacionais"

        # Center align
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_introduction_slide(self, prs):
        """Add the introduction slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Introdução"
        intro_text = self._summarize_text(
            self.content['sections']['introduction'], 200)
        content.text = intro_text if intro_text else "Análise de throughput em sistemas computacionais"

    def _add_section_slide(self, prs, title, content):
        """Add a section slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]

        title_shape.text = title
        content_shape.text = content

    def _add_formulas_slide(self, prs):
        """Add the formulas slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Fórmulas Principais"
        formulas_text = "\n\n".join(
            self.content['formulas'][:5])  # Show top 5 formulas
        content.text = formulas_text

    def _summarize_text(self, text, max_chars=200):
        """Summarize text to a maximum number of characters"""
        if not text:
            return ""

        # Split into sentences
        sentences = text.split('.')
        summary = ""

        for sentence in sentences:
            if len(summary) + len(sentence) <= max_chars:
                summary += sentence.strip() + ". "
            else:
                break

        return summary.strip()


if __name__ == "__main__":
    test_content = {
        'sections': {
            'abstract': 'This is a test abstract.',
            'introduction': 'This is a test introduction.',
            'methodology': 'This is a test methodology.',
            'results': 'These are test results.',
            'conclusion': 'This is a test conclusion.'
        },
        'formulas': ['T = λ/μ', 'R = λW'],
        'num_pages': 5
    }
    tool = PPTXCreator(content=test_content,
                       output_path="output/test_presentation.pptx")
    print(tool.run())
