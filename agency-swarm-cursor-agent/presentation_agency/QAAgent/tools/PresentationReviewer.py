from agency_swarm.tools import BaseTool
from pydantic import Field
from pptx import Presentation
import os
from dotenv import load_dotenv

load_dotenv()


class PresentationReviewer(BaseTool):
    """
    A tool for reviewing PowerPoint presentations, checking content quality, formatting, and language.
    Looks for presentations in the output/ directory.
    """

    presentation_filename: str = Field(
        ...,
        description="Name of the PowerPoint file to review (e.g., 'presentation.pptx')"
    )

    def run(self):
        """
        Reviews the PowerPoint presentation and returns a detailed analysis.
        The presentation should be in the output/ directory.
        """
        presentation_path = os.path.join("output", self.presentation_filename)

        try:
            prs = Presentation(presentation_path)
            review_results = []

            # Review each slide
            for slide_number, slide in enumerate(prs.slides, 1):
                slide_review = {
                    'slide_number': slide_number,
                    'content': [],
                    'formatting': [],
                    'suggestions': []
                }

                # Check title
                if slide.shapes.title:
                    title_text = slide.shapes.title.text
                    if len(title_text) > 50:
                        slide_review['formatting'].append(
                            f"Title is too long ({len(title_text)} characters)")
                else:
                    slide_review['formatting'].append(
                        "Slide is missing a title")

                # Check content
                text_content = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"

                # Content checks
                if len(text_content) < 10:
                    slide_review['content'].append(
                        "Slide has very little content")
                elif len(text_content) > 200:
                    slide_review['content'].append(
                        "Slide may have too much content")

                # Add the slide review to results
                if slide_review['content'] or slide_review['formatting'] or slide_review['suggestions']:
                    review_results.append(slide_review)

            # Format the review results as a string
            review_text = "Presentation Review Results:\n\n"
            for result in review_results:
                review_text += f"Slide {result['slide_number']}:\n"
                if result['content']:
                    review_text += "Content Issues:\n- " + \
                        "\n- ".join(result['content']) + "\n"
                if result['formatting']:
                    review_text += "Formatting Issues:\n- " + \
                        "\n- ".join(result['formatting']) + "\n"
                if result['suggestions']:
                    review_text += "Suggestions:\n- " + \
                        "\n- ".join(result['suggestions']) + "\n"
                review_text += "\n"

            return review_text

        except Exception as e:
            return f"Error reviewing presentation: {str(e)}"


if __name__ == "__main__":
    # Test the tool with a sample presentation
    tool = PresentationReviewer(presentation_filename="test_presentation.pptx")
    print(tool.run())
