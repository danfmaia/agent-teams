from agency_swarm.tools import BaseTool
from pydantic import Field, BaseModel
from typing import Dict, Optional, List, Union, Annotated
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import io
import os
from dotenv import load_dotenv

load_dotenv()


class GraphData(BaseModel):
    x: List[Union[int, float]]
    y: List[Union[int, float]]
    title: Optional[str] = ""
    xlabel: Optional[str] = ""
    ylabel: Optional[str] = ""


class SlideContent(BaseModel):
    title: str
    content: str
    graph_data: Optional[GraphData] = None


class PPTXCreator(BaseTool):
    """
    A tool for creating and formatting PowerPoint presentations with slides, text, formulas, and graphs.
    Output files are saved in the output/ directory.
    """

    content: Dict[int, SlideContent] = Field(
        description="Dictionary containing slide content. Each key is a slide number, and value is a SlideContent object with title, content, and optional graph_data"
    )

    output_filename: str = Field(
        description="The name of the PowerPoint file to be saved (e.g., 'presentation.pptx')"
    )

    def _create_title_slide(self, prs, title, subtitle=""):
        """Creates a title slide with optional subtitle."""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        if subtitle and slide.placeholders.get(1):
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle

        return slide

    def _create_content_slide(self, prs, title, content):
        """Creates a content slide with better formatting."""
        try:
            # Try section header layout first
            slide = prs.slides.add_slide(prs.slide_layouts[2])
        except:
            try:
                # Fallback to title and content layout
                slide = prs.slides.add_slide(prs.slide_layouts[1])
            except:
                # Last resort: blank slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add title if the slide has a title placeholder
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
        else:
            # Create a text box for title if no title placeholder
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = title

        # Add content
        try:
            # Try to use content placeholder
            content_shape = slide.placeholders[1]
        except:
            # Create a text box if no placeholder available
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            content_shape = slide.shapes.add_textbox(left, top, width, height)

        tf = content_shape.text_frame
        tf.text = content

        return slide

    def _add_graph(self, slide, graph_data: GraphData):
        """Adds a graph to the slide with better positioning and error handling."""
        try:
            plt.figure(figsize=(8, 4))

            plt.plot(graph_data.x, graph_data.y)
            plt.title(graph_data.title)
            plt.xlabel(graph_data.xlabel)
            plt.ylabel(graph_data.ylabel)

            # Save the graph to a byte stream
            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight')
            img_stream.seek(0)

            # Add the graph to the slide with better positioning
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            slide.shapes.add_picture(img_stream, left, top, width=width)
            plt.close()

            return None
        except Exception as e:
            plt.close()
            return f"Error adding graph: {str(e)}"

    def run(self):
        """
        Creates a PowerPoint presentation with the specified content, including text, formulas, and graphs.
        Saves the file in the output/ directory.
        """
        try:
            # Ensure output directory exists
            output_dir = "output"
            os.makedirs(output_dir, exist_ok=True)

            # Create full output path
            output_path = os.path.join(output_dir, self.output_filename)

            # Create presentation
            prs = Presentation()

            # Process each slide
            errors = []
            content_dict = dict(self.content)  # Convert to regular dictionary
            for slide_num, slide_content in content_dict.items():
                try:
                    # Create appropriate slide based on position
                    if slide_num == 1:
                        slide = self._create_title_slide(
                            prs, slide_content.title, slide_content.content)
                    else:
                        slide = self._create_content_slide(
                            prs, slide_content.title, slide_content.content)

                    # Add graph if present
                    if slide_content.graph_data:
                        error = self._add_graph(
                            slide, slide_content.graph_data)
                        if error:
                            errors.append(f"Slide {slide_num}: {error}")

                except Exception as e:
                    errors.append(
                        f"Error creating slide {slide_num}: {str(e)}")

            # Save the presentation
            prs.save(output_path)

            # Return result with any errors
            result = f"Presentation created successfully at {output_path}"
            if errors:
                result += "\n\nWarnings/Errors:\n" + "\n".join(errors)

            return result

        except Exception as e:
            return f"Critical error creating presentation: {str(e)}"


if __name__ == "__main__":
    # Test the tool with sample content
    sample_content = {
        1: SlideContent(
            title='Sample Presentation',
            content='Created with PPTXCreator'
        ),
        2: SlideContent(
            title='Sample Slide',
            content='This is a test slide with a graph',
            graph_data=GraphData(
                x=[1, 2, 3, 4, 5],
                y=[2, 4, 6, 8, 10],
                title='Sample Graph',
                xlabel='X axis',
                ylabel='Y axis'
            )
        )
    }

    tool = PPTXCreator(
        content=sample_content,
        output_filename="test_presentation.pptx"
    )
    print(tool.run())
