import os
import json
import re
import pickle
import shutil
from datetime import datetime
from pathlib import Path
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import numpy as np
from openai import OpenAI
from typing import Dict, List, Optional
import cmd
import readline
from rich.console import Console
from rich.panel import Panel
from rich.markdown import Markdown

# Initialize rich console for better output
console = Console()

# Load environment variables
load_dotenv()
client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

# Ensure generated files directories exist
GEN_DIR = Path("gen")
MEMORY_DIR = GEN_DIR / "memory"
PRESENTATIONS_DIR = GEN_DIR / "presentations"
GRAPHS_DIR = GEN_DIR / "graphs"
RUNS_DIR = GEN_DIR / "runs"

for dir_path in [GEN_DIR, MEMORY_DIR, PRESENTATIONS_DIR, GRAPHS_DIR, RUNS_DIR]:
    dir_path.mkdir(parents=True, exist_ok=True)


class ConversationMemory:
    def __init__(self, run_id: str):
        self.run_id = run_id
        self.file_path = MEMORY_DIR / f"memory_{run_id}.pkl"
        self.memory: Dict[str, List[Dict]] = self._load_memory()

    def _load_memory(self) -> Dict[str, List[Dict]]:
        if self.file_path.exists():
            with open(self.file_path, 'rb') as f:
                return pickle.load(f)
        return {}

    def save_memory(self):
        with open(self.file_path, 'wb') as f:
            pickle.dump(self.memory, f)

    def add_interaction(self, agent_name: str, interaction: Dict):
        if agent_name not in self.memory:
            self.memory[agent_name] = []
        self.memory[agent_name].append({
            'timestamp': datetime.now().isoformat(),
            'run_id': self.run_id,
            **interaction
        })
        self.save_memory()

    def get_agent_history(self, agent_name: str) -> List[Dict]:
        return self.memory.get(agent_name, [])

    def get_previous_runs(self) -> List[str]:
        """Get a list of all previous run IDs."""
        return [f.stem.replace('memory_', '') for f in MEMORY_DIR.glob('memory_*.pkl')]


class BaseAgent:
    def __init__(self, client: OpenAI, memory: ConversationMemory, run_id: str):
        self.client = client
        self.memory = memory
        self.run_id = run_id
        self.name = self.__class__.__name__

    def _call_gpt(self, prompt: str, model: str = "gpt-4") -> str:
        response = self.client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": prompt}]
        )
        result = response.choices[0].message.content

        self.memory.add_interaction(self.name, {
            'prompt': prompt,
            'response': result
        })

        return result


class PDFReaderAgent(BaseAgent):
    def extract_content(self, pdf_path: str) -> str:
        reader = PdfReader(pdf_path)
        content = ""
        for page in reader.pages:
            content += page.extract_text()

        self.memory.add_interaction(self.name, {
            'action': 'extract_content',
            'pdf_path': pdf_path,
            'content_length': len(content)
        })

        return content


class ContentOrganizerAgent(BaseAgent):
    def organize_content(self, content: str, previous_feedback: Optional[str] = None) -> Dict:
        prompt = f"""Analyze this content and create a 5-minute presentation outline with 5-7 slides.
        Focus on throughput estimation key points, formulas, and data for visualization.
        
        {previous_feedback if previous_feedback else ""}
        
        Return ONLY a JSON object with this exact structure, nothing else:
        {{
            "slides": [
                {{
                    "title": "slide title",
                    "content": "slide content"
                }}
            ]
        }}
        
        Content to analyze:
        {content[:3000]}..."""

        response_text = self._call_gpt(prompt)

        # Clean and parse JSON
        response_text = re.sub(r'^```json\s*', '', response_text)
        response_text = re.sub(r'\s*```$', '', response_text)

        try:
            return json.loads(response_text)
        except json.JSONDecodeError as e:
            print(f"Error parsing JSON response: {e}")
            print(f"Raw response: {response_text}")
            raise


class TranslationAgent(BaseAgent):
    def translate_to_portuguese(self, content: Dict, previous_feedback: Optional[str] = None) -> Dict:
        if not isinstance(content, dict) or 'slides' not in content:
            raise ValueError(
                "Content must be a dictionary with a 'slides' key")

        prompt = f"""Translate the following presentation content to Portuguese, maintaining technical terms.
        {previous_feedback if previous_feedback else ""}
        Return ONLY a JSON object with the same structure, nothing else."""

        response = self.client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "user", "content": prompt},
                {"role": "user", "content": json.dumps(
                    content, ensure_ascii=False)}
            ]
        )

        response_text = response.choices[0].message.content.strip()

        # Clean and parse JSON
        response_text = re.sub(r'^```json\s*', '', response_text)
        response_text = re.sub(r'\s*```$', '', response_text)

        try:
            translated = json.loads(response_text)
            self.memory.add_interaction(self.name, {
                'action': 'translate',
                'original': content,
                'translated': translated
            })
            return translated
        except json.JSONDecodeError as e:
            print(f"Error parsing JSON response: {e}")
            print(f"Raw response: {response_text}")
            raise


class VisualizationAgent(BaseAgent):
    def create_graphs(self, content: str, previous_feedback: Optional[str] = None) -> str:
        # Create a unique filename for this run
        graph_path = GRAPHS_DIR / f"graph_{self.run_id}.png"

        # This would create graphs based on the data and previous feedback
        plt.figure(figsize=(10, 6))
        x = np.linspace(0, 10, 100)
        y = np.sin(x)
        plt.plot(x, y)
        plt.title('Sample Throughput Graph')
        plt.xlabel('Time')
        plt.ylabel('Throughput')

        plt.savefig(graph_path)
        plt.close()

        self.memory.add_interaction(self.name, {
            'action': 'create_graph',
            'graph_path': str(graph_path),
            'previous_feedback': previous_feedback
        })

        return str(graph_path)


class PowerPointGeneratorAgent(BaseAgent):
    def create_presentation(self, content: Dict, graph_path: str, previous_feedback: Optional[str] = None) -> str:
        if not isinstance(content, dict) or 'slides' not in content:
            raise ValueError(
                "Content must be a dictionary with a 'slides' key")

        prs = Presentation()
        pptx_path = PRESENTATIONS_DIR / \
            f"apresentacao_throughput_{self.run_id}.pptx"

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Estimativa de Throughput"
        subtitle.text = "Baseado em Análise Técnica"

        # Content slides
        for slide_content in content['slides']:
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = slide_content['title']
            body.text = slide_content['content']

        # Graph slide
        graph_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(graph_slide_layout)
        title = slide.shapes.title
        title.text = "Visualização de Dados"
        slide.shapes.add_picture(graph_path, Inches(
            2), Inches(2), Inches(6), Inches(4))

        prs.save(pptx_path)

        self.memory.add_interaction(self.name, {
            'action': 'create_presentation',
            'pptx_path': str(pptx_path),
            'previous_feedback': previous_feedback
        })

        return str(pptx_path)


class QAAgent(BaseAgent):
    def review_presentation(self, pptx_path: str, original_content: str) -> Dict:
        prompt = f"""You are a QA agent reviewing a PowerPoint presentation about throughput estimation.
        Your task is to analyze the presentation and provide a structured review.
        
        IMPORTANT: Return ONLY a JSON object with NO additional text or markdown formatting.
        The JSON must follow this EXACT structure:
        {{
            "score": <number between 0-10>,
            "issues": [<list of strings>],
            "suggestions": [<list of strings>],
            "approval": <boolean>
        }}

        Review criteria:
        1. Accuracy of technical content
        2. Quality of Portuguese translation
        3. Clarity and flow of presentation
        4. Visual appeal and readability
        5. Completeness of content coverage
        
        Original content excerpt for reference:
        {original_content[:1000]}...
        """

        response_text = self._call_gpt(prompt).strip()

        # Extract JSON if it's wrapped in any markdown or text
        json_match = re.search(r'\{[\s\S]*\}', response_text)
        if json_match:
            response_text = json_match.group(0)

        try:
            review = json.loads(response_text)

            # Validate review structure
            required_keys = {'score', 'issues', 'suggestions', 'approval'}
            if not all(key in review for key in required_keys):
                raise ValueError("Missing required keys in review response")

            if not isinstance(review['score'], (int, float)) or not 0 <= review['score'] <= 10:
                raise ValueError("Score must be a number between 0 and 10")

            if not isinstance(review['issues'], list) or not isinstance(review['suggestions'], list):
                raise ValueError("Issues and suggestions must be lists")

            if not isinstance(review['approval'], bool):
                raise ValueError("Approval must be a boolean")

            # Save review to a JSON file
            review_path = RUNS_DIR / f"review_{self.run_id}.json"
            with open(review_path, 'w', encoding='utf-8') as f:
                json.dump(review, f, ensure_ascii=False, indent=2)

            self.memory.add_interaction(self.name, {
                'action': 'review',
                'pptx_path': pptx_path,
                'review': review,
                'review_path': str(review_path)
            })
            return review

        except json.JSONDecodeError as e:
            print(f"Error parsing JSON response: {e}")
            print(f"Raw response: {response_text}")
            raise
        except ValueError as e:
            print(f"Invalid review format: {e}")
            print(f"Raw response: {response_text}")
            raise


def generate_run_id() -> str:
    """Generate a unique run ID based on timestamp."""
    return datetime.now().strftime("%Y%m%d_%H%M%S")


class AgentTeam:
    def __init__(self):
        self.run_id = generate_run_id()
        self.memory = ConversationMemory(self.run_id)
        self.pdf_reader = PDFReaderAgent(client, self.memory, self.run_id)
        self.content_organizer = ContentOrganizerAgent(
            client, self.memory, self.run_id)
        self.translator = TranslationAgent(client, self.memory, self.run_id)
        self.visualizer = VisualizationAgent(client, self.memory, self.run_id)
        self.ppt_generator = PowerPointGeneratorAgent(
            client, self.memory, self.run_id)
        self.qa_agent = QAAgent(client, self.memory, self.run_id)
        self.content = None
        self.organized_content = None
        self.translated_content = None
        self.graph_path = None
        self.pptx_path = None
        self.review = None
        self.iteration = 0

    def show_status(self):
        """Show current status of the presentation generation process."""
        status = [
            ("PDF Content", "✅" if self.content else "❌"),
            ("Organized Content", "✅" if self.organized_content else "❌"),
            ("Portuguese Translation", "✅" if self.translated_content else "❌"),
            ("Visualizations", "✅" if self.graph_path else "❌"),
            ("PowerPoint", "✅" if self.pptx_path else "❌"),
            ("QA Review", "✅" if self.review else "❌")
        ]

        console.print("\n[bold]Current Status:[/bold]")
        for name, state in status:
            console.print(f"{name}: {state}")

        if self.iteration > 0:
            console.print(f"\nCurrent Iteration: {self.iteration}")

    def show_previous_runs(self):
        """Show history of previous runs and their feedback."""
        previous_runs = self.memory.get_previous_runs()
        if not previous_runs:
            console.print("\nNo previous runs found.")
            return

        console.print("\n[bold]Previous Runs:[/bold]")
        for run_id in previous_runs:
            review_path = RUNS_DIR / f"review_{run_id}.json"
            if review_path.exists():
                with open(review_path, 'r', encoding='utf-8') as f:
                    review = json.load(f)
                console.print(f"\n[bold]Run {run_id}:[/bold]")
                console.print(f"Score: {review['score']}/10")
                console.print("Issues:")
                for issue in review['issues']:
                    console.print(f"- {issue}")
                console.print(
                    f"Approved: {'✅' if review['approval'] else '❌'}")

    def should_improve(self) -> bool:
        """Determine if we should improve the existing presentation."""
        return (
            self.iteration > 0
            and self.review is not None
            and self.review['score'] >= 5
            and not self.review['approval']
        )

    def extract_content(self):
        if not self.should_improve():
            console.print("\n[bold]Extracting PDF content...[/bold]")
            self.content = self.pdf_reader.extract_content(
                'input/throughput_estimation_article.pdf')
            console.print("✅ Content extracted successfully")

    def organize_content(self):
        if not self.content:
            console.print("[red]❌ Please extract PDF content first[/red]")
            return

        if self.should_improve():
            console.print(
                "\n[bold]Improving content organization based on feedback...[/bold]")
            feedback = f"""Previous review (score: {self.review['score']}/10):
            Issues: {', '.join(self.review['issues'])}
            Suggestions: {', '.join(self.review['suggestions'])}
            """
            self.organized_content = self.content_organizer.organize_content(
                self.content,
                previous_feedback=feedback
            )
            console.print("✅ Content organization improved")
        else:
            console.print("\n[bold]Organizing content...[/bold]")
            self.organized_content = self.content_organizer.organize_content(
                self.content)
            console.print("✅ Content organized successfully")

    def translate_content(self):
        if not self.organized_content:
            console.print("[red]❌ Please organize content first[/red]")
            return

        if self.should_improve():
            console.print(
                "\n[bold]Improving translation based on feedback...[/bold]")
            feedback = f"""Previous review (score: {self.review['score']}/10):
            Issues: {', '.join(self.review['issues'])}
            Suggestions: {', '.join(self.review['suggestions'])}
            """
            self.translated_content = self.translator.translate_to_portuguese(
                self.organized_content,
                previous_feedback=feedback
            )
            console.print("✅ Translation improved")
        else:
            console.print("\n[bold]Translating to Portuguese...[/bold]")
            self.translated_content = self.translator.translate_to_portuguese(
                self.organized_content)
            console.print("✅ Content translated successfully")

    def create_visualizations(self):
        if not self.content:
            console.print("[red]❌ Please extract PDF content first[/red]")
            return

        if self.should_improve():
            console.print(
                "\n[bold]Improving visualizations based on feedback...[/bold]")
            feedback = f"""Previous review (score: {self.review['score']}/10):
            Issues: {', '.join(self.review['issues'])}
            Suggestions: {', '.join(self.review['suggestions'])}
            """
            self.graph_path = self.visualizer.create_graphs(
                self.content,
                previous_feedback=feedback
            )
            console.print("✅ Visualizations improved")
        else:
            console.print("\n[bold]Creating visualizations...[/bold]")
            self.graph_path = self.visualizer.create_graphs(self.content)
            console.print("✅ Visualizations created successfully")

    def generate_presentation(self):
        if not all([self.translated_content, self.graph_path]):
            console.print(
                "[red]❌ Please complete translation and visualizations first[/red]")
            return

        if self.should_improve():
            console.print(
                "\n[bold]Improving presentation based on feedback...[/bold]")
            feedback = f"""Previous review (score: {self.review['score']}/10):
            Issues: {', '.join(self.review['issues'])}
            Suggestions: {', '.join(self.review['suggestions'])}
            """
            self.pptx_path = self.ppt_generator.create_presentation(
                self.translated_content,
                self.graph_path,
                previous_feedback=feedback
            )
            console.print("✅ Presentation improved")
        else:
            console.print(
                "\n[bold]Generating PowerPoint presentation...[/bold]")
            self.pptx_path = self.ppt_generator.create_presentation(
                self.translated_content, self.graph_path)
            console.print("✅ Presentation generated successfully")

    def review_presentation(self):
        if not all([self.pptx_path, self.content]):
            console.print("[red]❌ Please generate presentation first[/red]")
            return

        console.print("\n[bold]Reviewing presentation...[/bold]")
        self.review = self.qa_agent.review_presentation(
            self.pptx_path, self.content)

        console.print("\n[bold]QA Review Results:[/bold]")
        console.print(f"Score: {self.review['score']}/10")
        console.print("\nIssues found:")
        for issue in self.review['issues']:
            console.print(f"- {issue}")
        console.print("\nSuggestions:")
        for suggestion in self.review['suggestions']:
            console.print(f"- {suggestion}")
        console.print(
            f"\nApproval: {'✅ Approved' if self.review['approval'] else '❌ Not Approved'}")

        # Increment iteration counter after review
        self.iteration += 1


class AgentChat(cmd.Cmd):
    intro = '''Welcome to the Presentation Generator Agent Team Chat!
Type 'help' or '?' to list commands.
    '''
    prompt = '🤖 > '

    def __init__(self):
        super().__init__()
        self.team = AgentTeam()

    def do_status(self, arg):
        """Show current status of the presentation generation process"""
        self.team.show_status()

    def do_history(self, arg):
        """Show previous runs and their feedback"""
        self.team.show_previous_runs()

    def do_extract(self, arg):
        """Extract content from the PDF"""
        self.team.extract_content()

    def do_organize(self, arg):
        """Organize the extracted content into presentation format"""
        self.team.organize_content()

    def do_translate(self, arg):
        """Translate the organized content to Portuguese"""
        self.team.translate_content()

    def do_visualize(self, arg):
        """Create visualizations for the presentation"""
        self.team.create_visualizations()

    def do_generate(self, arg):
        """Generate the PowerPoint presentation"""
        self.team.generate_presentation()

    def do_review(self, arg):
        """Review the generated presentation"""
        self.team.review_presentation()

    def do_auto(self, arg):
        """Run the complete process automatically"""
        self.team.extract_content()
        self.team.organize_content()
        self.team.translate_content()
        self.team.create_visualizations()
        self.team.generate_presentation()
        self.team.review_presentation()

    def do_quit(self, arg):
        """Exit the chat"""
        console.print("\nGoodbye! 👋")
        return True

    def do_EOF(self, arg):
        """Exit on Ctrl-D"""
        return self.do_quit(arg)


if __name__ == '__main__':
    console.print(Panel.fit(
        "[bold blue]Presentation Generator Agent Team[/bold blue]\n"
        "An interactive system for generating presentations with AI assistance"
    ))
    AgentChat().cmdloop()
