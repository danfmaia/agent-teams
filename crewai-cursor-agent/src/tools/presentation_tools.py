from langchain.tools import BaseTool
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib.pyplot as plt
import io
import base64
from typing import Optional, Dict, List


class PresentationTools:
    @staticmethod
    def create_slide() -> BaseTool:
        def add_slide(content: Dict[str, str], pptx_path: str = "output/presentation.pptx") -> str:
            """Create or add a slide to the presentation"""
            try:
                try:
                    prs = Presentation(pptx_path)
                except:
                    prs = Presentation()

                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                content_shape = slide.shapes.placeholders[1]

                title.text = content.get("title", "")
                content_shape.text = content.get("content", "")

                prs.save(pptx_path)
                return f"Slide added successfully to {pptx_path}"
            except Exception as e:
                return f"Error creating slide: {str(e)}"

        return BaseTool(
            name="Create Slide",
            func=add_slide,
            description="Creates a new slide in the PowerPoint presentation"
        )

    @staticmethod
    def add_graph() -> BaseTool:
        def create_graph(data: Dict[str, List], title: str, pptx_path: str = "output/presentation.pptx") -> str:
            """Add a graph to the presentation"""
            try:
                try:
                    prs = Presentation(pptx_path)
                except:
                    prs = Presentation()

                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # Add title
                title_shape = slide.shapes.title
                title_shape.text = title

                # Create chart data
                chart_data = CategoryChartData()
                chart_data.categories = data["categories"]
                chart_data.add_series("Series 1", data["values"])

                # Add chart
                x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
                ).chart

                prs.save(pptx_path)
                return f"Graph added successfully to {pptx_path}"
            except Exception as e:
                return f"Error adding graph: {str(e)}"

        return BaseTool(
            name="Add Graph",
            func=create_graph,
            description="Adds a graph to the PowerPoint presentation"
        )

    @staticmethod
    def review_presentation() -> BaseTool:
        def review(pptx_path: str = "output/presentation.pptx") -> str:
            """Review the presentation and return a summary"""
            try:
                prs = Presentation(pptx_path)
                summary = []

                for i, slide in enumerate(prs.slides, 1):
                    slide_content = f"Slide {i}:\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_content += f"- {shape.text}\n"
                    summary.append(slide_content)

                return "\n".join(summary)
            except Exception as e:
                return f"Error reviewing presentation: {str(e)}"

        return BaseTool(
            name="Review Presentation",
            func=review,
            description="Reviews the PowerPoint presentation and provides a summary"
        )
